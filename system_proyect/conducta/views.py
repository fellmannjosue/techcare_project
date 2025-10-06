import io, os, json
from django.shortcuts import render, redirect, get_object_or_404
from django.http import JsonResponse, HttpResponse, HttpResponseForbidden
from django.contrib.auth.decorators import login_required
from django.template.loader import render_to_string
from django.contrib import messages
from django.db import connections
from django.db.models import Count
from django.utils import timezone
from django.conf import settings

# PDF y PowerPoint
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import mm
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .forms import (
    ReporteInformativoForm,
    ReporteConductualForm,
    ProgressReportForm
)
from .models import (
    IncisoConductual,
    ReporteInformativo,
    ReporteConductual,
    ProgressReport,
    MateriaDocenteBilingue,
    MateriaDocenteColegio
)
from tickets.models import Ticket


# ---------------------------
# FUNCIONES AUXILIARES
# ---------------------------
def obtener_alumnos_bilingue():
    query = """
    SELECT d.PersonaID, 
           ISNULL(d.Nombre1,'') + 
           CASE WHEN d.Nombre2 IS NOT NULL AND d.Nombre2 <> '' THEN ' ' + d.Nombre2 ELSE '' END + ' ' +
           ISNULL(d.Apellido1,'') + 
           CASE WHEN d.Apellido2 IS NOT NULL AND d.Apellido2 <> '' THEN ' ' + d.Apellido2 ELSE '' END AS NombreCompl,
           da.Descripcion AS AreaDesc, c.CrsoNumero, c.GrupoNumero
      FROM dbo.tblPrsDtosGen AS d
      JOIN dbo.tblPrsTipo           AS t  ON d.PersonaID = t.PersonaID
      JOIN dbo.tblEdcArea           AS a  ON t.IngrEgrID  = a.IngrEgrID
      JOIN dbo.tblEdcEjecCrso       AS ec ON a.AreaID     = ec.AreaID
      JOIN dbo.tblEdcCrso           AS c  ON ec.CrsoID    = c.CrsoID
      JOIN dbo.tblEdcDescripAreaEdc AS da ON a.DescrAreaEdcID = da.DescrAreaEdcID
     WHERE d.Alum = 1 AND DATEPART(yy, c.FechaInicio) = DATEPART(yyyy, GETDATE())
       AND da.Descripcion IN (N'PrimariaBL', N'ColegioBL', N'PreescolarBL')
       AND ec.Desertor = 0 AND ec.TrasladoPer = 0
     ORDER BY da.Descripcion, c.CrsoNumero, c.GrupoNumero, NombreCompl
    """
    alumnos = []
    try:
        with connections['padres_sqlserver'].cursor() as cursor:
            cursor.execute(query)
            for pid, nombre, area, crso, grupo in cursor.fetchall():
                label = nombre.strip()
                grado = f"{area} {crso}-{grupo}"
                alumnos.append({'id': str(pid), 'label': label, 'grado': grado})
    except Exception as e:
        print(">>> ERROR SQL BILINGUE:", e)
    return alumnos


def obtener_alumnos_colegio():
    query = """
    SELECT TOP (100) PERCENT d.PersonaID,
           ISNULL(d.Nombre1, '') + 
           CASE WHEN d.Nombre2 IS NOT NULL AND d.Nombre2 <> '' THEN ' ' + d.Nombre2 ELSE '' END + ' ' +
           ISNULL(d.Apellido1, '') + 
           CASE WHEN d.Apellido2 IS NOT NULL AND d.Apellido2 <> '' THEN ' ' + d.Apellido2 ELSE '' END AS NombreCompl,
           da.Descripcion AS AreaDesc, c.CrsoNumero, c.GrupoNumero
      FROM dbo.tblPrsDtosGen AS d
     INNER JOIN dbo.tblPrsTipo           AS t  ON d.PersonaID = t.PersonaID
     INNER JOIN dbo.tblEdcArea           AS a  ON t.IngrEgrID  = a.IngrEgrID
     INNER JOIN dbo.tblEdcEjecCrso       AS ec ON a.AreaID     = ec.AreaID
     INNER JOIN dbo.tblEdcCrso           AS c  ON ec.CrsoID    = c.CrsoID
     INNER JOIN dbo.tblEdcDescripAreaEdc AS da ON a.DescrAreaEdcID = da.DescrAreaEdcID
     WHERE (d.Alum = 1) 
       AND (DATEPART(yy, c.FechaInicio) = DATEPART(yyyy, GETDATE()))
       AND (da.Descripcion IN (N'Colegio', N'Bachillerato')) 
       AND (ec.Desertor = 0) AND (ec.TrasladoPer = 0)
     ORDER BY AreaDesc, c.CrsoNumero, c.GrupoNumero, NombreCompl
    """
    alumnos = []
    try:
        with connections['padres_sqlserver'].cursor() as cursor:
            cursor.execute(query)
            for pid, nombre, area, crso, grupo in cursor.fetchall():
                label = nombre.strip()
                grado = f"{area} {crso}-{grupo}"
                alumnos.append({'id': str(pid), 'label': label, 'grado': grado})
    except Exception as e:
        print(">>> ERROR SQL COLEGIO:", e)
    return alumnos


def get_materia_docente_choices(area):
    if area == 'bilingue':
        return [
            (str(md.pk), f"{md.materia} — {md.docente}")
            for md in MateriaDocenteBilingue.objects.filter(activo=True)
        ]
    else:
        return [
            (str(md.pk), f"{md.materia} — {md.docente}")
            for md in MateriaDocenteColegio.objects.filter(activo=True)
        ]

# ---------------------------
# DASHBOARDS Y FORMULARIOS
# ---------------------------
@login_required
def dashboard_maestro(request):
    user = request.user
    area = None
    if user.groups.filter(name='maestros_bilingue').exists():
        area = 'bilingue'
    elif user.groups.filter(name='maestros_colegio').exists():
        area = 'colegio'
    return render(request, 'conducta/dashboard_maestros.html', {'area': area})

# ------------ REPORTE INFORMATIVO  ------------
@login_required
def reporte_informativo_bilingue(request):
    area = 'bilingue'
    students = obtener_alumnos_bilingue()  
    materia_docente_choices = get_materia_docente_choices(area)
    fecha = timezone.now().strftime("%Y-%m-%dT%H:%M")

    if request.method == 'POST':
        alumno_id = request.POST.get('alumno')
        grado = request.POST.get('grado')
        materia_docente_id = request.POST.get('materia_docente')
        comentario = request.POST.get('comentario', "")
        alumno_obj = next((a for a in students if a['id'] == alumno_id), None)
        alumno_label = alumno_obj['label'] if alumno_obj else ""
        materia = docente = ""
        if materia_docente_id:
            md_obj = MateriaDocenteBilingue.objects.filter(pk=materia_docente_id).first()
            if md_obj:
                materia = md_obj.materia
                docente = md_obj.docente
        ReporteInformativo.objects.create(
            usuario=request.user,
            area=area,
            alumno_id=alumno_id,
            alumno_nombre=alumno_label,
            grado=grado,
            materia=materia,
            docente=docente,
            comentario=comentario
        )
        messages.success(request, "¡Reporte registrado correctamente!")
        return redirect('reporte_informativo_bilingue')

    return render(request, 'conducta/form_informativo.html', {
        'students': students,    
        'materia_docente_choices': materia_docente_choices,
        'fecha': fecha,
        'area': area,
    })

@login_required
def reporte_informativo_colegio(request):
    area = 'colegio'
    students = obtener_alumnos_colegio()
    materia_docente_choices = get_materia_docente_choices(area)
    fecha = timezone.now().strftime("%Y-%m-%dT%H:%M")

    if request.method == 'POST':
        alumno_id = request.POST.get('alumno')
        grado = request.POST.get('grado')
        materia_docente_id = request.POST.get('materia_docente')
        comentario = request.POST.get('comentario', "")
        alumno_obj = next((a for a in students if a['id'] == alumno_id), None)
        alumno_label = alumno_obj['label'] if alumno_obj else ""
        materia = docente = ""
        if materia_docente_id:
            md_obj = MateriaDocenteColegio.objects.filter(pk=materia_docente_id).first()
            if md_obj:
                materia = md_obj.materia
                docente = md_obj.docente
        ReporteInformativo.objects.create(
            usuario=request.user,
            area=area,
            alumno_id=alumno_id,
            alumno_nombre=alumno_label,
            grado=grado,
            materia=materia,
            docente=docente,
            comentario=comentario
        )
        messages.success(request, "¡Reporte registrado correctamente!")
        return redirect('reporte_informativo_colegio')

    return render(request, 'conducta/form_informativo.html', {
        'students': students,
        'materia_docente_choices': materia_docente_choices,
        'fecha': fecha,
        'area': area,
    })

# ------------ REPORTE CONDUCTUAL  ------------
@login_required
def reporte_conductual_bilingue(request):
    area = 'bilingue'
    students = obtener_alumnos_bilingue()
    materia_docente_choices = get_materia_docente_choices(area)
    fecha = timezone.now().strftime("%Y-%m-%d")
    # Obtiene incisos por severidad
    incisos_leve = IncisoConductual.objects.filter(activo=True, tipo='leve').order_by('descripcion')
    incisos_grave = IncisoConductual.objects.filter(activo=True, tipo='grave').order_by('descripcion')
    incisos_muygrave = IncisoConductual.objects.filter(activo=True, tipo='muy_grave').order_by('descripcion')

    if request.method == 'POST':
        alumno_id = request.POST.get('alumno')
        grado = request.POST.get('grado')
        materia_docente_id = request.POST.get('materia_docente')
        fecha_val = request.POST.get('fecha')
        comentario = request.POST.get('comentario', "")

        # Obtén las listas seleccionadas si el checkbox está activo
        ids_leve = request.POST.getlist('inciso_leve[]') if request.POST.get('chk_leve') else []
        ids_grave = request.POST.getlist('inciso_grave[]') if request.POST.get('chk_grave') else []
        ids_muygrave = request.POST.getlist('inciso_muygrave[]') if request.POST.get('chk_muygrave') else []

        alumno_obj = next((a for a in students if a['id'] == alumno_id), None)
        alumno_label = alumno_obj['label'] if alumno_obj else ""
        materia = docente = ""
        if materia_docente_id:
            md_obj = MateriaDocenteBilingue.objects.filter(pk=materia_docente_id).first()
            if md_obj:
                materia = md_obj.materia
                docente = md_obj.docente

        # Crea el reporte sin los ManyToMany
        reporte = ReporteConductual.objects.create(
            usuario=request.user,
            area=area,
            alumno_id=alumno_id,
            alumno_nombre=alumno_label,
            grado=grado,
            materia=materia,
            docente=docente,
            fecha=fecha_val,
            comentario=comentario,
        )
        # Asocia los incisos ManyToMany
        if ids_leve:
            reporte.incisos_leve.set(ids_leve)
        if ids_grave:
            reporte.incisos_grave.set(ids_grave)
        if ids_muygrave:
            reporte.incisos_muygrave.set(ids_muygrave)

        messages.success(request, "¡Reporte conductual registrado correctamente!")
        return redirect('reporte_conductual_bilingue')

    return render(request, 'conducta/form_conductual.html', {
        'students': students,
        'materia_docente_choices': materia_docente_choices,
        'fecha': fecha,
        'area': area,
        'incisos_leve': incisos_leve,
        'incisos_grave': incisos_grave,
        'incisos_muygrave': incisos_muygrave,
    })

@login_required
def reporte_conductual_colegio(request):
    area = 'colegio'
    students = obtener_alumnos_colegio()
    materia_docente_choices = get_materia_docente_choices(area)
    fecha = timezone.now().strftime("%Y-%m-%d")
    incisos_leve = IncisoConductual.objects.filter(activo=True, tipo='leve').order_by('descripcion')
    incisos_grave = IncisoConductual.objects.filter(activo=True, tipo='grave').order_by('descripcion')
    incisos_muygrave = IncisoConductual.objects.filter(activo=True, tipo='muy_grave').order_by('descripcion')

    if request.method == 'POST':
        alumno_id = request.POST.get('alumno')
        grado = request.POST.get('grado')
        materia_docente_id = request.POST.get('materia_docente')
        fecha_val = request.POST.get('fecha')
        comentario = request.POST.get('comentario', "")

        ids_leve = request.POST.getlist('inciso_leve[]') if request.POST.get('chk_leve') else []
        ids_grave = request.POST.getlist('inciso_grave[]') if request.POST.get('chk_grave') else []
        ids_muygrave = request.POST.getlist('inciso_muygrave[]') if request.POST.get('chk_muygrave') else []

        alumno_obj = next((a for a in students if a['id'] == alumno_id), None)
        alumno_label = alumno_obj['label'] if alumno_obj else ""
        materia = docente = ""
        if materia_docente_id:
            md_obj = MateriaDocenteColegio.objects.filter(pk=materia_docente_id).first()
            if md_obj:
                materia = md_obj.materia
                docente = md_obj.docente

        reporte = ReporteConductual.objects.create(
            usuario=request.user,
            area=area,
            alumno_id=alumno_id,
            alumno_nombre=alumno_label,
            grado=grado,
            materia=materia,
            docente=docente,
            fecha=fecha_val,
            comentario=comentario,
        )
        if ids_leve:
            reporte.incisos_leve.set(ids_leve)
        if ids_grave:
            reporte.incisos_grave.set(ids_grave)
        if ids_muygrave:
            reporte.incisos_muygrave.set(ids_muygrave)

        messages.success(request, "¡Reporte conductual registrado correctamente!")
        return redirect('reporte_conductual_colegio')

    return render(request, 'conducta/form_conductual.html', {
        'students': students,
        'materia_docente_choices': materia_docente_choices,
        'fecha': fecha,
        'area': area,
        'incisos_leve': incisos_leve,
        'incisos_grave': incisos_grave,
        'incisos_muygrave': incisos_muygrave,
    })

#-------------- PROGRESS REPORT -----------------
@login_required
def progress_report_bilingue(request):
    MATERIAS_PRIMARIA = [
        "Math", "Phonics", "Spelling", "Reading", "Language",
        "Science", "Español", "CCSS", "Asociadas"
    ]
    MATERIAS_COLEGIO = [
        "Math", "Spelling", "Reading", "Language", "Science",
        "Español", "CCSS", "Cívica", "Asociadas"
    ]

    students = obtener_alumnos_bilingue()
    alumnos_choices = [(s['id'], s['label']) for s in students]

    materias = MATERIAS_PRIMARIA
    grado = ""
    alumno_id = ""
    if request.method == 'POST':
        form = ProgressReportForm(alumnos_choices=alumnos_choices, data=request.POST)
        alumno_id = request.POST.get('alumno')
        grado = request.POST.get('grado', '')
        if alumno_id:
            alumno_obj = next((s for s in students if s['id'] == alumno_id), None)
            if alumno_obj:
                grado = alumno_obj['grado']
        if grado and ("Primaria" in grado or "Preescolar" in grado):
            materias = MATERIAS_PRIMARIA
        else:
            materias = MATERIAS_COLEGIO

        if form.is_valid():
            # --- PROCESAR MATERIAS ---
            materias_list = []
            for materia in materias:
                if materia == "Asociadas":
                    asignaciones = request.POST.getlist('asignacion_Asociadas[]')
                    comentarios = request.POST.getlist('comentario_Asociadas[]')
                    for idx, asignacion in enumerate(asignaciones):
                        comentario = comentarios[idx] if idx < len(comentarios) else ""
                        materias_list.append({
                            'materia': 'Asociadas',
                            'asignacion': asignacion,
                            'comentario': comentario,
                        })
                else:
                    asignacion = request.POST.get(f"asignacion_{materia}", "")
                    comentario = request.POST.get(f"comentario_{materia}", "")
                    materias_list.append({
                        'materia': materia,
                        'asignacion': asignacion,
                        'comentario': comentario,
                    })
            # --- DATOS GENERALES ---
            semana_inicio = form.cleaned_data.get('semana_inicio')
            semana_fin = form.cleaned_data.get('semana_fin')
            comentario_general = form.cleaned_data.get('comentario_general', "")
            alumno_id = form.cleaned_data.get('alumno')
            alumno_obj = next((s for s in students if s['id'] == alumno_id), None)
            alumno_label = alumno_obj['label'] if alumno_obj else ""
            grado = form.cleaned_data.get('grado')

            # --- CREAR REPORTE ---
            ProgressReport.objects.create(
                usuario=request.user,
                alumno_id=alumno_id,
                alumno_nombre=alumno_label,
                grado=grado,
                semana_inicio=semana_inicio,
                semana_fin=semana_fin,
                comentario_general=comentario_general,
                materias_json=materias_list  # <-- SE GUARDA EL JSON AQUÍ
            )
            messages.success(request, "¡Progress report registrado correctamente!")
            return redirect('progress_report_bilingue')
    else:
        form = ProgressReportForm(alumnos_choices=alumnos_choices)
        materias = MATERIAS_PRIMARIA

    return render(request, 'conducta/form_progress.html', {
        'form': form,
        'materias': materias,
        'students': students,
        'grado': grado,
    })


#-------------- HISTORIAL MAESTROS -----------------
@login_required
def historial_maestro_bilingue(request):
    usuario = request.user
    usuario_actual = usuario.get_full_name()
    reportes_informativo = ReporteInformativo.objects.filter(usuario=usuario, area='bilingue').order_by('-fecha')
    reportes_conductual = ReporteConductual.objects.filter(usuario=usuario, area='bilingue').order_by('-fecha')

    reportes_progress = ProgressReport.objects.all().order_by('-fecha')

    tickets_usuario = Ticket.objects.filter(email=usuario.email).order_by('-created_at')

    return render(request, 'conducta/historial_maestro.html', {
        'reportes_informativo': reportes_informativo,
        'reportes_conductual': reportes_conductual,
        'reportes_progress': reportes_progress,
        'tickets_usuario': tickets_usuario,
        'area': 'bilingue',
    })


@login_required
def historial_maestro_colegio(request):
    usuario = request.user
    reportes_informativo = ReporteInformativo.objects.filter(usuario=usuario, area='colegio').order_by('-fecha')
    reportes_conductual = ReporteConductual.objects.filter(usuario=usuario, area='colegio').order_by('-fecha')

    tickets_usuario = Ticket.objects.filter(email=usuario.email).order_by('-created_at')

    return render(request, 'conducta/historial_maestro.html', {
        'reportes_informativo': reportes_informativo,
        'reportes_conductual': reportes_conductual,
        'reportes_progress': [],   # vacía para mantener compatibilidad
        'tickets_usuario': tickets_usuario,      # <- AGREGA AQUÍ
        'area': 'colegio',
    })

# ----------- EDITOR DE REPORTES ( SOLO COORDINADOR) -----------
def es_coordinador(user):
    return user.is_staff or user.groups.filter(name__icontains="coordinador").exists()


# Cambia estos según tu lógica o roles
COORDINADORES_BL = ["Mrs. Osorto", "Miss Alcerro", "Miss Angela"]
COORDINADORES_COLEGIO = ["Profe. Licona", "Profe. Felipe", "Profe. Gabriela"]

# ────────────────
# EDITAR CONDUCTUAL
# ────────────────

@login_required
def editar_reporte_conductual(request, pk):
    reporte = get_object_or_404(ReporteConductual, pk=pk)

    # Elige coordinadores según área
    if reporte.area == "bilingue":
        coordinadores = COORDINADORES_BL
    else:
        coordinadores = COORDINADORES_COLEGIO

    if request.method == "POST":
        # Guardar los campos editables del coordinador
        comentario_coordinador = request.POST.get("comentario_coordinador", "")
        estado = request.POST.get("estado", "enviado")
        coordinador_firma = request.POST.get("coordinador_firma", "")
        # Actualiza solo los campos de coordinación
        reporte.comentario_coordinador = comentario_coordinador
        reporte.estado = estado
        reporte.coordinador_firma = coordinador_firma
        reporte.save()
        messages.success(request, "Reporte conductual actualizado correctamente.")
        return redirect('dashboard_coordinador', area=reporte.area)
    # Renderizar formulario
    return render(request, "conducta/editor_conductual.html", {
        "reporte": reporte,
        "coordinadores": coordinadores,
    })

# ────────────────
# EDITAR INFORMATIVO
# ────────────────

@login_required
def editar_reporte_informativo(request, pk):
    reporte = get_object_or_404(ReporteInformativo, pk=pk)

    # Elige coordinadores según área
    if reporte.area == "bilingue":
        coordinadores = COORDINADORES_BL
    else:
        coordinadores = COORDINADORES_COLEGIO

    if request.method == "POST":
        comentario_coordinador = request.POST.get("comentario_coordinador", "")
        estado = request.POST.get("estado", "enviado")
        coordinador_firma = request.POST.get("coordinador_firma", "")
        reporte.comentario_coordinador = comentario_coordinador
        reporte.estado = estado
        reporte.coordinador_firma = coordinador_firma
        reporte.save()
        messages.success(request, "Reporte informativo actualizado correctamente.")
        return redirect('dashboard_coordinador', area=reporte.area)
    # Renderizar formulario
    return render(request, "conducta/editor_informativo.html", {
        "reporte": reporte,
        "coordinadores": coordinadores,
    })

# ────────────────
# EDITAR PROGRESS
# ────────────────

@login_required
def editar_progress_report(request, pk):
    import json
    from django.contrib import messages
    from django.shortcuts import render, redirect, get_object_or_404

    COORDINADORES_BL = ["Mrs. Osorto", "Miss Alcerro", "Miss Angela"]
    coordinadores = COORDINADORES_BL

    reporte = get_object_or_404(ProgressReport, pk=pk)
    es_coord = es_coordinador(request.user)
    usuario_actual = request.user.get_full_name() or request.user.username  # Usa el nombre completo si está

    # --- Decodificar materias ---
    materias = reporte.materias_json
    if isinstance(materias, str):
        try:
            materias = json.loads(materias)
        except Exception:
            materias = []
    if materias is None:
        materias = []

    # --- Determina para cada materia si el usuario la puede editar ---
    for mat in materias:
        if es_coord or not mat.get("docente") or mat.get("docente", "") == usuario_actual:
            mat["editable"] = True
        else:
            mat["editable"] = False

    if request.method == 'POST':
        nuevas_materias = []
        asociadas_indices_existentes = []
        # 1️⃣ Procesa materias normales y existentes (NO 'Asociadas')
        for mat in materias:
            materia = mat.get('materia', '')
            if materia == 'Asociadas':
                # Guarda índices para procesarlos después
                asociadas_indices_existentes.append(mat)
                continue
            if mat["editable"]:
                asignacion = request.POST.get(f'asignacion_{materia}', mat.get('asignacion', ''))
                comentario = request.POST.get(f'comentario_{materia}', mat.get('comentario', ''))
                mat['asignacion'] = asignacion
                mat['comentario'] = comentario
                # SOLO asigna el docente si está vacío
                if not mat.get('docente'):
                    mat['docente'] = usuario_actual
            nuevas_materias.append(mat)

        # 2️⃣ Procesa TODAS las filas de 'Asociadas'
        asignaciones_asociadas = request.POST.getlist('asignacion_Asociadas[]')
        comentarios_asociadas = request.POST.getlist('comentario_Asociadas[]')
        # Los nombres de los docentes asociados pueden venir en POST (si pones un input), si no, asigna usuario_actual a todos los nuevos
        # Si quieres que sólo el que edita quede como docente, déjalos así:
        for i in range(len(asignaciones_asociadas)):
            # Si tienes filas de 'Asociadas' ya existentes, usa el docente de esas filas; si no, pon el usuario actual
            docente_valor = usuario_actual
            if i < len(asociadas_indices_existentes) and asociadas_indices_existentes[i].get('docente'):
                docente_valor = asociadas_indices_existentes[i]['docente']
            nuevas_materias.append({
                'materia': 'Asociadas',
                'asignacion': asignaciones_asociadas[i],
                'comentario': comentarios_asociadas[i],
                'docente': docente_valor,
                'editable': True,  # Esto es sólo para el template
            })

        # 3️⃣ Actualiza los campos generales y guarda
        reporte.materias_json = json.dumps(nuevas_materias)
        reporte.comentario_general = request.POST.get('comentario_general', reporte.comentario_general)
        if es_coord:
            reporte.coordinador_firma = request.POST.get('coordinador_firma', reporte.coordinador_firma)
            reporte.estado = request.POST.get('estado', reporte.estado)
            reporte.comentario_coordinador = request.POST.get('comentario_coordinador', reporte.comentario_coordinador)

        reporte.save()
        messages.success(request, "¡Reporte actualizado correctamente!")
        if es_coord:
            return redirect('dashboard_coordinador', area='bilingue')
        else:
            return redirect('historial_maestro_bilingue')

    return render(request, 'conducta/editor_progress.html', {
        'reporte': reporte,
        'materias': materias,
        'es_coordinador': es_coord,
        'usuario_actual': usuario_actual,
        'coordinadores': coordinadores,
    })

# -----------  DESCARGA EN PDF  -----------
def draw_paragraph(pdf, text, x, y, max_width, font="Helvetica", font_size=10, bold=False, italic=False, leading=13):
    """
    Dibuja un párrafo con saltos de línea y justificación manual en ReportLab.
    """
    from reportlab.pdfbase.pdfmetrics import stringWidth
    fontname = font
    if bold and italic:
        fontname = "Helvetica-BoldOblique"
    elif bold:
        fontname = "Helvetica-Bold"
    elif italic:
        fontname = "Helvetica-Oblique"
    pdf.setFont(fontname, font_size)
    lines = []
    for raw_line in text.split('\n'):
        line = ""
        for word in raw_line.split():
            test_line = f"{line} {word}".strip()
            if stringWidth(test_line, fontname, font_size) < max_width:
                line = test_line
            else:
                lines.append(line)
                line = word
        lines.append(line)
    for l in lines:
        pdf.drawString(x, y, l)
        y -= leading
    return y


@login_required
def descargar_pdf_informativo(request, pk):
    from .models import ReporteInformativo
    reporte = get_object_or_404(ReporteInformativo, pk=pk)
    buf = io.BytesIO()
    w, h = letter
    pdf = canvas.Canvas(buf, pagesize=letter)
    pdf.setTitle("reporte_informativo.pdf")

    # Logo centrado
    width_logo = 35 * mm
    height_logo = 35 * mm
    x_logo = (w - width_logo) / 2
    logo_path = os.path.join(settings.STATIC_ROOT, "conducta/img/ana-transformed.png")
    if os.path.exists(logo_path):
        pdf.drawImage(
            logo_path,
            x=x_logo,
            y=h-40*mm,
            width=width_logo,
            height=height_logo,
            mask='auto'
        )
    # 🔥 Más margen superior
    y_actual = h - 52*mm

    # Título y datos
    pdf.setFont("Helvetica-Bold", 16)
    pdf.drawCentredString(w/2, y_actual, "Nuevo Amanecer School" if reporte.area == "bilingue" else "C.E.M.N.G Nuevo Amanecer")
    y_actual -= 14*mm
    pdf.setFont("Helvetica", 13)
    pdf.drawCentredString(w/2, y_actual, "Reporte informativo")
    y_actual -= 14*mm

    # Datos alumno
    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(32*mm, y_actual, "Nombre:")
    pdf.setFont("Helvetica", 10)
    pdf.drawString(60*mm, y_actual, reporte.alumno_nombre)

    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(110*mm, y_actual, "Grado:")
    pdf.setFont("Helvetica", 10)
    pdf.drawString(130*mm, y_actual, reporte.grado)
    y_actual -= 9*mm

    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(32*mm, y_actual, "Docente:")
    pdf.setFont("Helvetica", 10)
    pdf.drawString(65*mm, y_actual, reporte.docente or '')

    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(90*mm, y_actual, "Fecha:")
    pdf.setFont("Helvetica", 11)
    pdf.drawString(110*mm, y_actual, reporte.fecha.strftime('%d/%m/%Y') if reporte.fecha else '')
    y_actual -= 16*mm

    # Comentario del docente
    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(32*mm, y_actual, "Comentario del Docente:")
    y_actual -= 8*mm
    pdf.setFont("Helvetica-Oblique", 10)
    y_actual = draw_paragraph(
        pdf, reporte.comentario or "-", x=38*mm, y=y_actual, max_width=w-65*mm,
        font="Helvetica", font_size=10, italic=True, leading=12
    )
    y_actual -= 10*mm

    # Comentario del coordinador si existe
    if hasattr(reporte, 'comentario_coordinador') and reporte.comentario_coordinador:
        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(32*mm, y_actual, "Comentario del Coordinador:")
        y_actual -= 8*mm
        pdf.setFont("Helvetica-Oblique", 10)
        y_actual = draw_paragraph(
            pdf, reporte.comentario_coordinador, x=38*mm, y=y_actual, max_width=w-65*mm,
            font="Helvetica", font_size=10, italic=True, leading=12
        )
        y_actual -= 8*mm

    # ====== FIRMAS DOCENTE Y COORDINADOR ======
    y_firma = 40 * mm  # o el valor que uses para que siempre quede abajo
    x_firma_doc = 38 * mm
    x_firma_coord = 110 * mm
    largo_firma = 65 * mm  # Puedes usar el mismo largo para ambos

    # Firma docente
    pdf.setFont("Helvetica-Bold", 10)
    pdf.drawString(x_firma_doc, y_firma + 5*mm, "Firma:")
    pdf.setFont("Helvetica", 10)
    pdf.drawString(x_firma_doc + 35*mm, y_firma + 5*mm, f"{reporte.docente or ''}")
    pdf.setStrokeColor(colors.black)
    pdf.setLineWidth(0.7)
    pdf.line(x_firma_doc, y_firma, x_firma_doc + largo_firma, y_firma)

    # Firma coordinador
    pdf.setFont("Helvetica-Bold", 10)
    pdf.drawString(x_firma_coord, y_firma + 5*mm, "Firma:")
    pdf.setFont("Helvetica", 10)
    pdf.drawString(x_firma_coord + 38*mm, y_firma + 5*mm, f"{getattr(reporte, 'coordinador_firma', '') or ''}")
    pdf.setStrokeColor(colors.black)
    pdf.setLineWidth(0.7)
    pdf.line(x_firma_coord, y_firma, x_firma_coord + largo_firma, y_firma)


    pdf.save()
    buf.seek(0)
    return HttpResponse(buf, content_type="application/pdf")


@login_required
def descargar_pdf_conductual(request, pk):
    from .models import ReporteConductual
    reporte = get_object_or_404(ReporteConductual, pk=pk)
    buf = io.BytesIO()
    w, h = letter
    pdf = canvas.Canvas(buf, pagesize=letter)
    pdf.setTitle("reporte_conductual.pdf")

    width_logo = 35 * mm
    height_logo = 35 * mm
    x_logo = (w - width_logo) / 2
    logo_path = os.path.join(settings.STATIC_ROOT, "conducta/img/ana-transformed.png")
    if os.path.exists(logo_path):
        pdf.drawImage(
            logo_path,
            x=x_logo,
            y=h-40*mm,
            width=width_logo,
            height=height_logo,
            mask='auto'
        )
    y_actual = h - 52*mm

    # Título
    pdf.setFont("Helvetica-Bold", 16)
    titulo = "Nuevo Amanecer School" if getattr(reporte, "area", None) == "bilingue" else "C.E.M.N.G Nuevo Amanecer"
    pdf.drawCentredString(w/2, y_actual, titulo)
    y_actual -= 14*mm
    pdf.setFont("Helvetica", 13)
    pdf.drawCentredString(w/2, y_actual, "Reporte conductual")
    y_actual -= 14*mm

    # -----------------------------
    # Nombre (primera línea)
    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(32*mm, y_actual, "Nombre:")
    pdf.setFont("Helvetica", 11)
    pdf.drawString(60*mm, y_actual, reporte.alumno_nombre)

    # Grado y Fecha (debajo)
    y_actual -= 8*mm
    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(32*mm, y_actual, "Grado:")
    pdf.setFont("Helvetica", 11)
    pdf.drawString(60*mm, y_actual, reporte.grado)
    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(110*mm, y_actual, "Fecha:")
    pdf.setFont("Helvetica", 11)
    pdf.drawString(130*mm, y_actual, reporte.fecha.strftime('%d/%m/%Y') if reporte.fecha else '')
    
    # Docente (debajo)
    y_actual -= 8*mm
    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(32*mm, y_actual, "Docente:")
    pdf.setFont("Helvetica", 11)
    pdf.drawString(60*mm, y_actual, reporte.docente or "")
    y_actual -= 12*mm

    # ----------- INCISOS -----------
    maxw = w - 65*mm
    for tipo, label in [("incisos_leve", "Leve"), ("incisos_grave", "Grave"), ("incisos_muygrave", "Muy Grave")]:
        incisos = getattr(reporte, tipo).all()
        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(32*mm, y_actual, f"Incisos {label}:")
        y_actual -= 8*mm
        if incisos:
            for i in incisos:
                y_actual = draw_paragraph(
                    pdf, i.descripcion, x=38*mm, y=y_actual, max_width=maxw,
                    font="Helvetica", font_size=10, bold=True, italic=True, leading=12
                )
        else:
            pdf.setFont("Helvetica-Oblique", 10)
            pdf.drawString(38*mm, y_actual, "-")
            y_actual -= 8*mm
        y_actual -= 6*mm

    # ----------- COMENTARIOS -----------
    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(32*mm, y_actual, "Comentario del Docente:")
    y_actual -= 8*mm
    y_actual = draw_paragraph(
        pdf, reporte.comentario or "-", x=38*mm, y=y_actual, max_width=maxw,
        font="Helvetica", font_size=10, italic=True, leading=12
    )
    y_actual -= 10*mm

    if hasattr(reporte, 'comentario_coordinador') and reporte.comentario_coordinador:
        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(32*mm, y_actual, "Comentario del Coordinador:")
        y_actual -= 8*mm
        y_actual = draw_paragraph(
            pdf, reporte.comentario_coordinador, x=38*mm, y=y_actual, max_width=maxw,
            font="Helvetica", font_size=10, italic=True, leading=12
        )
        y_actual -= 8*mm

    # ======================
    #  FIRMA (siempre abajo)
    # ======================

    y_firma = 40 * mm   # Puedes subir/bajar este valor según tu hoja.
    x_firma_doc = 38 * mm
    x_firma_coord = 110 * mm
    largo_firma = 65 * mm  # Puedes hacer la línea un poco más larga si lo deseas

    # ------ FIRMA DOCENTE ------
    pdf.setFont("Helvetica-Bold", 10)
    pdf.drawString(x_firma_doc, y_firma + 5*mm, "Firma:")
    pdf.setFont("Helvetica", 10)
    pdf.drawString(x_firma_doc + 22*mm, y_firma + 5*mm, f"{reporte.docente or ''}")
    # Línea exactamente debajo
    pdf.setStrokeColor(colors.black)
    pdf.setLineWidth(0.7)
    pdf.line(x_firma_doc, y_firma, x_firma_doc + largo_firma, y_firma)

    # ------ FIRMA COORDINADOR ------
    pdf.setFont("Helvetica-Bold", 10)
    pdf.drawString(x_firma_coord, y_firma + 5*mm, "Firma:")
    pdf.setFont("Helvetica", 10)
    pdf.drawString(x_firma_coord + 22*mm, y_firma + 5*mm, f"{getattr(reporte, 'coordinador_firma', '') or ''}")
    pdf.setStrokeColor(colors.black)
    pdf.setLineWidth(0.7)
    pdf.line(x_firma_coord, y_firma, x_firma_coord + largo_firma, y_firma)
    
    pdf.save()
    buf.seek(0)
    return HttpResponse(buf, content_type="application/pdf")


@login_required
def descargar_pdf_progress(request, pk):
    reporte = get_object_or_404(ProgressReport, pk=pk)
    materias = reporte.materias_json or []

    # Ruta correcta
    fondo_path = '/home/admin2/techcare_project/system_proyect/conducta/static/conducta/img/plantilla.jpg'

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Slide en blanco
    slide = prs.slides.add_slide(slide_layout)

    # Fondo completo
    slide.shapes.add_picture(fondo_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # -- TÍTULO "Progress Report"
    titulo = slide.shapes.add_textbox(Inches(0.7), Inches(0.1), Inches(5), Inches(0.7)).text_frame
    p = titulo.paragraphs[0]
    p.text = "Progress Report"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = "French Script MT"  # Usa Arial Black o una fuente bien visible
    p.alignment = PP_ALIGN.LEFT

    # -- Name: <cursiva>   (en la misma línea)
    nombre_textbox = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(7.5), Inches(0.7))
    tf = nombre_textbox.text_frame
    tf.clear()

    # "Name: " en negrita
    run1 = tf.paragraphs[0].add_run()
    run1.text = "Name: "
    run1.font.bold = True
    run1.font.size = Pt(18)
    run1.font.name = "Arial (Cuerpo) "

    # El nombre en cursiva/script (usa "Brush Script MT", o prueba con "Segoe Script" o similar)
    run2 = tf.paragraphs[0].add_run()
    run2.text = reporte.alumno_nombre
    run2.font.size = Pt(28)
    run2.font.bold = True
    run2.font.name = "Edwardian Script ITC"  # Cambia aquí si no tienes esta fuente

    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # -- Weeks: <rango> (a la derecha)
    semanas_textbox = slide.shapes.add_textbox(Inches(6), Inches(0.7), Inches(4), Inches(0.5))
    semanas_tf = semanas_textbox.text_frame
    semanas_tf.clear()
    semanas_tf.word_wrap = True
    p = semanas_tf.paragraphs[0]
    p.text = f"Weeks: {reporte.semana_inicio.strftime('%b %d')} - {reporte.semana_fin.strftime('%b %d, %Y')}"
    p.font.italic = True
    p.font.size = Pt(18)
    p.font.name = "Bahnschrift Light SemiCondensed"
    p.alignment = PP_ALIGN.RIGHT

    # -- Grado
    grado_textbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(5), Inches(0.5))
    grado_tf = grado_textbox.text_frame
    grado_tf.clear()
    p = grado_tf.paragraphs[0]
    p.text = f"Grade: {reporte.grado}"
    p.font.size = Pt(18)
    p.font.name = "Bahnschrift SemiBold SemiConden"
    p.alignment = PP_ALIGN.LEFT

    # -- TABLA de materias (fondos blancos, cabecera en negrita)
    num_rows = len(materias) + 1
    num_cols = 3
    left = Inches(3.15)
    top = Inches(1.57)
    width = Inches(6.496063)
    height = Inches(5.11811)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Cabecera
    headers = ["Materia", "Asignación", "Comentario/Observación"]
    for idx, title in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = title
        # --- ESTILO DE CABECERA ---
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(220, 220, 220)  # Gris claro
        p = cell.text_frame.paragraphs[0]
        p.font.bold = False
        p.font.size = Pt(18)
        p.font.name = "Arial (Cuerpo)"
        p.font.color.rgb = RGBColor(0, 0, 0)  # NEGRO
        p.alignment = PP_ALIGN.CENTER

        # Fondo blanco y bordes (no azul)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(236, 236, 233)
        for border in cell._tc.xpath('.//a:tcPr/a:ln'):
            border.set('w', '12700')  # Ajusta borde si deseas más grueso

    # Filas de materias
    for row, mat in enumerate(materias, start=1):
        # Materia en negrita
        cell_materia = table.cell(row, 0)
        cell_materia.text = mat.get("materia", "")
        cell_materia.text_frame.paragraphs[0].font.bold = False
        cell_materia.text_frame.paragraphs[0].font.size = Pt(17)
        cell_materia.fill.solid()
        cell_materia.fill.fore_color.rgb = RGBColor(236, 236, 233)
        # Asignación
        cell_asignacion = table.cell(row, 1)
        cell_asignacion.text = mat.get("asignacion", "")
        cell_asignacion.text_frame.paragraphs[0].font.size = Pt(17)
        cell_asignacion.fill.solid()
        cell_asignacion.fill.fore_color.rgb = RGBColor(236, 236, 233)
        # Comentario
        cell_comentario = table.cell(row, 2)
        cell_comentario.text = mat.get("comentario", "")
        cell_comentario.text_frame.paragraphs[0].font.size = Pt(17)
        cell_comentario.fill.solid()
        cell_comentario.fill.fore_color.rgb = RGBColor(236, 236, 233)
        # Quitar bordes azules si aparecen (PowerPoint suele ponerlos por defecto)
        for cell in [cell_materia, cell_asignacion, cell_comentario]:
            for border in cell._tc.xpath('.//a:tcPr/a:ln'):
                border.set('w', '12700')  # Ajusta borde si deseas más grueso

    # DESCARGA el archivo
    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    filename = f'progress_{reporte.alumno_nombre.replace(" ", "_")}.pptx'
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    prs.save(response)
    return response


@login_required
def descargar_pdf_conductual_3_strikes(request, pk):
    from .models import ReporteConductual
    reporte = get_object_or_404(ReporteConductual, pk=pk)
    reportes = list(ReporteConductual.objects.filter(
        area=reporte.area,
        alumno_id=reporte.alumno_id
    ).order_by('fecha')[:3])

    if len(reportes) < 3:
        return HttpResponse("Este alumno aún no tiene 3 reportes conductuales.", content_type="text/plain")

    buf = io.BytesIO()
    w, h = letter
    pdf = canvas.Canvas(buf, pagesize=letter)
    pdf.setTitle("reporte_conductual_3_strikes.pdf")

    # Logo centrado
    width_logo = 35 * mm
    height_logo = 35 * mm
    x_logo = (w - width_logo) / 2
    logo_path = os.path.join(settings.STATIC_ROOT, "conducta/img/ana-transformed.png")
    if os.path.exists(logo_path):
        pdf.drawImage(
            logo_path,
            x=x_logo,
            y=h-40*mm,
            width=width_logo,
            height=height_logo,
            mask='auto'
        )
    y_actual = h - 47 * mm

    # Título
    pdf.setFont("Helvetica-Bold", 16)
    titulo = "Nuevo Amanecer School" if reporte.area == "bilingue" else "C.E.M.N.G Nuevo Amanecer"
    pdf.drawCentredString(w/2, y_actual, titulo)
    y_actual -= 12*mm
    pdf.setFont("Helvetica", 13)
    pdf.drawCentredString(w/2, y_actual, "Reporte conductual – 3 Strikes")
    y_actual -= 10*mm

    # --- CADA REPORTE ---
    for idx, rep in enumerate(reportes, 1):
        pdf.setFont("Helvetica-Bold", 12)
        pdf.drawString(32*mm, y_actual, f"Reporte #{idx} ({', '.join(['Leve' if rep.incisos_leve.exists() else '', 'Grave' if rep.incisos_grave.exists() else '', 'Muy Grave' if rep.incisos_muygrave.exists() else '']).replace(',,',',').strip(', ')})")
        y_actual -= 8*mm

        # Fila: Nombre/Grado/Docente/Fecha
        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(32*mm, y_actual, "Nombre:")
        pdf.setFont("Helvetica", 11)
        pdf.drawString(60*mm, y_actual, rep.alumno_nombre)

        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(120*mm, y_actual, "Grado:")
        pdf.setFont("Helvetica", 11)
        pdf.drawString(140*mm, y_actual, rep.grado)
        y_actual -= 6*mm

        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(32*mm, y_actual, "Docente/Materia:")
        pdf.setFont("Helvetica", 11)
        pdf.drawString(75*mm, y_actual, rep.docente)
        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(120*mm, y_actual, "Fecha:")
        pdf.setFont("Helvetica", 11)
        pdf.drawString(140*mm, y_actual, rep.fecha.strftime('%d/%m/%Y'))
        y_actual -= 7*mm

        # Incisos (párrafo)
        pdf.setFont("Helvetica-Bold", 10)
        pdf.drawString(32*mm, y_actual, "Inciso:")
        y_actual -= 6*mm

        incisos = [i.descripcion for i in rep.incisos_leve.all()] + [i.descripcion for i in rep.incisos_grave.all()] + [i.descripcion for i in rep.incisos_muygrave.all()]
        if incisos:
            text_obj = pdf.beginText(38*mm, y_actual)
            text_obj.setFont("Helvetica-BoldOblique", 10)
            wrap_len = 95  # Ajusta a gusto (número de caracteres por línea)
            for inciso in incisos:
                lines = [inciso[i:i+wrap_len] for i in range(0, len(inciso), wrap_len)]
                for ln in lines:
                    text_obj.textLine(ln)
                    y_actual -= 5*mm
            pdf.drawText(text_obj)
            y_actual = text_obj.getY() - 1*mm
        else:
            pdf.setFont("Helvetica", 10)
            pdf.drawString(38*mm, y_actual, "-")
            y_actual -= 5*mm

        # Descripción (Comentario)
        pdf.setFont("Helvetica-Bold", 10)
        pdf.drawString(32*mm, y_actual, "Descripción:")
        y_actual -= 5*mm
        if rep.comentario:
            text_obj = pdf.beginText(38*mm, y_actual)
            text_obj.setFont("Helvetica-Oblique", 10)
            for line in rep.comentario.split("\n"):
                text_obj.textLine(line)
                y_actual -= 5*mm
            pdf.drawText(text_obj)
            y_actual = text_obj.getY() - 1*mm
        else:
            pdf.setFont("Helvetica", 10)
            pdf.drawString(38*mm, y_actual, "-")
            y_actual -= 5*mm

        # Espacio entre reportes
        y_actual -= 4*mm
        pdf.setStrokeColor(colors.grey)
        pdf.line(30*mm, y_actual, w-30*mm, y_actual)
        y_actual -= 8*mm

    # FIRMAS PARA LLENAR EN FÍSICO
    y_firmas = 20*mm
    x_firma = [34*mm, 91*mm, 146*mm]
    if reporte.area == "colegio":
        etiquetas = ["Firma Padre de Familia", "Firma del Docente/Orientación", "Firma de Consejería"]
    else:
        etiquetas = ["Firma Padre de Familia", "Firma del Docente", "Firma del Coordinador"]

    for i in range(3):
        pdf.setStrokeColor(colors.black)
        pdf.line(x_firma[i], y_firmas, x_firma[i]+42*mm, y_firmas)
        pdf.setFont("Helvetica", 10)
        pdf.drawCentredString(x_firma[i]+21*mm, y_firmas-5*mm, etiquetas[i])

    pdf.save()
    buf.seek(0)
    return HttpResponse(buf, content_type="application/pdf")

#--------------  DASHBOARD COORDINADOR -----------------
@login_required
def dashboard_coordinador(request, area):
    """
    Dashboard coordinador: muestra todos los reportes de Informativo, Conductual y Progress
    según el área.
    """
    # Traer todos los reportes por área
    if area == 'bilingue':
        reportes_informativo = ReporteInformativo.objects.filter(area='bilingue')
        reportes_conductual = ReporteConductual.objects.filter(area='bilingue')
        reportes_progress = ProgressReport.objects.all()  # solo bilingüe tiene progress
    elif area == 'colegio':
        reportes_informativo = ReporteInformativo.objects.filter(area='colegio')
        reportes_conductual = ReporteConductual.objects.filter(area='colegio')
        reportes_progress = []  # colegio NO tiene progress
    else:
        return redirect('menu')  # área inválida

    # Calcula "strikes" (conteo de reportes conductuales por alumno)
    # Dict { alumno_id: count }
    strikes = {}
    qs = reportes_conductual.values('alumno_id').annotate(total=Count('id')).filter(total__gte=3)
    for row in qs:
        strikes[row['alumno_id']] = row['total']

    contexto = {
        'area': area,
        'reportes_informativo': reportes_informativo,
        'reportes_conductual': reportes_conductual,
        'reportes_progress': reportes_progress,
        'strikes': strikes,
    }
    return render(request, 'conducta/dashboard_coordinador.html', contexto)

    
@login_required
def historial_alumno_coordinador(request, alumno_id):
    """
    Devuelve el historial de reportes (informativo, conductual, progress)
    para un alumno específico. Se usa en el modal del dashboard coordinador.
    """
    informativos = ReporteInformativo.objects.filter(alumno_id=alumno_id).order_by('-fecha')
    conductuales = ReporteConductual.objects.filter(alumno_id=alumno_id).order_by('-fecha')
    progress = ProgressReport.objects.filter(alumno_id=alumno_id).order_by('-fecha')
    context = {
        'informativos': informativos,
        'conductuales': conductuales,
        'progress': progress,
    }
    html = render_to_string('conducta/lista_reportes.html', context)
    return HttpResponse(html)
